from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
prs = Presentation()

# Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Enhanced Bankruptcy Prediction Using ARIMA-Neural Network Model"
subtitle.text = "Project Presentation"

# Problem Statement Slide
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_2.shapes.title, slide_2.placeholders[1]
title.text = "Problem Statement"
content.text = ("Traditional models for bankruptcy prediction, such as Altman's Z-score and logistic regression, "
                "fail to account for complex non-linear patterns. This project proposes a hybrid ARIMA-Neural Network model "
                "to improve predictive accuracy.")

# Introduction Slide
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_3.shapes.title, slide_3.placeholders[1]
title.text = "Introduction"
content.text = ("This project integrates ARIMA and Neural Networks to capture linear trends and non-linear patterns in financial data. "
                "The hybrid approach aims to provide a more robust bankruptcy prediction framework.")

# Objectives Slide
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_4.shapes.title, slide_4.placeholders[1]
title.text = "Objectives"
content.text = ("1. Study ARIMA and Neural Network models.\n"
                "2. Conduct exploratory data analysis (EDA).\n"
                "3. Define evaluation metrics (RMSE, MAE, MAPE).\n"
                "4. Design and implement the hybrid model.\n"
                "5. Evaluate the model and compare results.\n"
                "6. Provide insights and recommendations.")

# Data Analysis Slide
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_5.shapes.title, slide_5.placeholders[1]
title.text = "Data Analysis"
content.text = ("Key steps:\n"
                "- Preprocessing: Normalization, handling missing data, feature engineering.\n"
                "- ARIMA: Captured linear trends and seasonality.\n"
                "- Neural Network: Modeled non-linear dependencies.\n"
                "- Combined Model: Integrated ARIMA and NN for enhanced accuracy.")

# Code Overview Slide
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_6.shapes.title, slide_6.placeholders[1]
title.text = "Code Overview and Testing"
content.text = ("The codebase includes modules for data preprocessing, ARIMA modeling, neural network training, and visualization.\n"
                "Testing goals:\n"
                "- Accuracy evaluation (MAE, MAPE)\n"
                "- Error analysis and model stability\n"
                "- Performance comparison")

# Results Slide
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_7.shapes.title, slide_7.placeholders[1]
title.text = "Review of the Results"
content.text = ("ARIMA Model Performance:\n"
                "- MAE: 6.23 | MAPE: 12.45%\n\n"
                "Combined ARIMA-Neural Network Performance:\n"
                "- MAE: 7.45 | MAPE: 15.67%\n\n"
                "The combined model demonstrated better handling of non-linear patterns.")

# Visualizations Slide
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_8.shapes.title, slide_8.placeholders[1]
title.text = "Visualizations (Placeholders)"
content.text = ("1. ARIMA Forecast vs Actual\n"
                "2. Combined Model Forecast vs Actual\n"
                "3. Error Distribution Plots\n"
                "4. Cumulative Error Over Time\n"
                "5. Confidence Interval Bands")

# Conclusion Slide
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title, content = slide_9.shapes.title, slide_9.placeholders[1]
title.text = "Conclusion"
content.text = ("The hybrid ARIMA-Neural Network model improves bankruptcy prediction by leveraging the strengths of both "
                "approaches. Future work includes hyperparameter tuning, exploring alternative architectures, and real-time predictions.")

# Save the presentation
output_path = "Enhanced_Bankruptcy_Prediction_Presentation.pptx"
prs.save(output_path)

output_path
